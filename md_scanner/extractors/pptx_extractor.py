# PPTX Extractor stub
# Uses python-pptx to extract text from .pptx files

def extract_pptx_text(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except ImportError:
        raise RuntimeError("python-pptx not installed")
    except Exception as e:
        return f"[PPTX extraction error: {e}]"

